"""
Inserts 2 SHAP slides into the existing PPT after Slide 7 (Fingerprinting Module).
Slide 8  → SHAP Global Feature Importance (bar chart)
Slide 9  → SHAP Per-Device Heatmap
Run: python add_shap_slides.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree
from pathlib import Path
import io

PPT_PATH  = Path("IoT_Fingerprinting_Presentation.pptx")
PLOTS_DIR = Path("plots")

# ── Palette (same as generate_ppt.py) ───────────────────────────────────────
C_BG      = RGBColor(0x0D, 0x0F, 0x14)
C_SURFACE = RGBColor(0x13, 0x16, 0x1E)
C_SURF2   = RGBColor(0x1A, 0x1E, 0x28)
C_BORDER  = RGBColor(0x1F, 0x25, 0x35)
C_CYAN    = RGBColor(0x00, 0xD4, 0xE8)
C_PURPLE  = RGBColor(0xA0, 0x5A, 0xFF)
C_GREEN   = RGBColor(0x00, 0xE5, 0x80)
C_YELLOW  = RGBColor(0xFF, 0xCD, 0x43)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT    = RGBColor(0xC9, 0xD1, 0xE0)
C_DIM     = RGBColor(0x6B, 0x78, 0x98)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation(PPT_PATH)
BLANK = prs.slide_layouts[6]


# ── Helpers ──────────────────────────────────────────────────────────────────

def bg(slide, color=C_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def txbox(slide, text, x, y, w, h,
          font_size=Pt(14), bold=False, color=C_TEXT,
          align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text        = text
    run.font.size   = font_size
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name   = "Calibri"
    return tb


def heading(slide, text, color=C_CYAN):
    txbox(slide, text, Inches(0.5), Inches(0.2), Inches(12.33), Inches(0.6),
          font_size=Pt(26), bold=True, color=color)


def divider(slide, y=Inches(0.88)):
    rect(slide, Inches(0.5), y, Inches(12.33), Pt(1.5), fill_color=C_BORDER)


def footer(slide):
    rect(slide, 0, H - Inches(0.35), W, Inches(0.35), fill_color=C_SURFACE)
    txbox(slide, "Md Ghufran Alam  |  NDU202400038  |  NIELIT Srinagar  |  M.Tech Cyber Forensics",
          Inches(0.3), H - Inches(0.33), Inches(10), Inches(0.32),
          font_size=Pt(9), color=C_DIM)
    txbox(slide, "IoT Device Fingerprinting Framework",
          Inches(10.3), H - Inches(0.33), Inches(2.7), Inches(0.32),
          font_size=Pt(9), color=C_DIM, align=PP_ALIGN.RIGHT)


def bullet_row(slide, items, y, colors=None):
    """Render a row of small stat-chips."""
    x = Inches(0.5)
    chip_w = Inches(12.33) / len(items)
    for i, (label, val) in enumerate(items):
        color = (colors[i] if colors else C_CYAN)
        rect(slide, x + i * chip_w + Inches(0.05), y,
             chip_w - Inches(0.1), Inches(0.55),
             fill_color=C_SURF2, line_color=color, line_width=Pt(1))
        txbox(slide, val,
              x + i * chip_w + Inches(0.05), y + Inches(0.0),
              chip_w - Inches(0.1), Inches(0.3),
              font_size=Pt(15), bold=True, color=color, align=PP_ALIGN.CENTER)
        txbox(slide, label,
              x + i * chip_w + Inches(0.05), y + Inches(0.28),
              chip_w - Inches(0.1), Inches(0.25),
              font_size=Pt(9), color=C_DIM, align=PP_ALIGN.CENTER)


# ── Insert slide at a given index ────────────────────────────────────────────
def insert_slide_at(prs, index, layout):
    """Add a blank slide and move it to `index`."""
    slide = prs.slides.add_slide(layout)
    # slides are stored in presentation XML — move new slide to desired position
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    new_slide_elem = slides_list[-1]        # just added → last
    xml_slides.remove(new_slide_elem)
    xml_slides.insert(index, new_slide_elem)
    return slide


# ═══════════════════════════════════════════════════════════════════════════
#  NEW SLIDE A — SHAP Global Feature Importance
# ═══════════════════════════════════════════════════════════════════════════
sl_a = insert_slide_at(prs, 7, BLANK)   # position 8 (0-indexed = 7)
bg(sl_a)
footer(sl_a)
heading(sl_a, "Model Explainability — SHAP Feature Importance")
divider(sl_a)

# Key insight chips
bullet_row(
    sl_a,
    [("#1 Feature", "tcp_ratio"), ("#2 Feature", "udp_ratio"),
     ("#3 Feature", "ack_ratio"), ("#4 Feature", "mean_dest_port"),
     ("#5 Feature", "is_mqtt")],
    Inches(1.0),
    colors=[C_CYAN, C_CYAN, C_PURPLE, C_PURPLE, C_GREEN]
)

# SHAP bar chart — large, left side
sl_a.shapes.add_picture(
    str(PLOTS_DIR / "shap_bar.png"),
    Inches(0.4), Inches(1.7), Inches(7.6), Inches(5.3)
)

# Right side: explanation text box
rect(sl_a, Inches(8.2), Inches(1.7), Inches(4.8), Inches(5.3),
     fill_color=C_SURFACE, line_color=C_CYAN, line_width=Pt(1.5))

txbox(sl_a, "What is SHAP?",
      Inches(8.35), Inches(1.8), Inches(4.5), Inches(0.38),
      font_size=Pt(13), bold=True, color=C_CYAN)

txbox(sl_a,
      "SHAP (SHapley Additive exPlanations) quantifies "
      "the contribution of each feature to every individual "
      "prediction — not just global averages.",
      Inches(8.35), Inches(2.25), Inches(4.5), Inches(0.75),
      font_size=Pt(11), color=C_TEXT)

txbox(sl_a, "Key Findings",
      Inches(8.35), Inches(3.1), Inches(4.5), Inches(0.35),
      font_size=Pt(13), bold=True, color=C_PURPLE)

findings = [
    ("tcp_ratio / udp_ratio", "Protocol choice is the strongest device fingerprint"),
    ("ack_ratio",             "TCP handshake pattern uniquely identifies camera vs sensor"),
    ("mean_dest_port",        "Port 443 = camera/TV, Port 1883 = thermostat, Port 5683 = bulb"),
    ("is_mqtt / is_https",    "Application-layer flags alone can identify 6 of 8 devices"),
    ("is_coap",               "CoAP exclusively used by smart_bulb — near-perfect separator"),
]
for i, (feat, meaning) in enumerate(findings):
    y = Inches(3.55) + i * Inches(0.58)
    txbox(sl_a, feat,
          Inches(8.35), y, Inches(4.5), Inches(0.25),
          font_size=Pt(10), bold=True, color=C_YELLOW)
    txbox(sl_a, meaning,
          Inches(8.35), y + Inches(0.24), Inches(4.5), Inches(0.28),
          font_size=Pt(10), color=C_DIM)

txbox(sl_a,
      "SHAP proves the model is not a black box — every prediction is fully explainable.",
      Inches(8.35), Inches(6.55), Inches(4.5), Inches(0.35),
      font_size=Pt(10), italic=True, color=C_DIM)


# ═══════════════════════════════════════════════════════════════════════════
#  NEW SLIDE B — SHAP Per-Device Heatmap
# ═══════════════════════════════════════════════════════════════════════════
sl_b = insert_slide_at(prs, 8, BLANK)   # position 9 (0-indexed = 8)
bg(sl_b)
footer(sl_b)
heading(sl_b, "Per-Device SHAP Heatmap — Feature Influence per Device Type", color=C_PURPLE)
divider(sl_b)

# Heatmap image — full width center
sl_b.shapes.add_picture(
    str(PLOTS_DIR / "shap_device_heatmap.png"),
    Inches(0.4), Inches(1.0), Inches(12.53), Inches(4.55)
)

# Bottom interpretation row
rect(sl_b, Inches(0.4), Inches(5.7), Inches(12.53), Inches(1.5),
     fill_color=C_SURFACE, line_color=C_BORDER, line_width=Pt(1))

txbox(sl_b, "How to Read:",
      Inches(0.6), Inches(5.78), Inches(2.0), Inches(0.32),
      font_size=Pt(11), bold=True, color=C_CYAN)

interpretations = [
    (C_CYAN,   "Smart Camera",     "tcp_ratio + is_https are bright → camera identified by HTTPS streaming"),
    (C_PURPLE, "Thermostat",       "is_mqtt + mean_dest_port (1883) dominate → MQTT periodic updates"),
    (C_GREEN,  "Smart Bulb",       "is_coap completely dominates → CoAP port 5683 is a unique signature"),
    (C_YELLOW, "Motion Sensor",    "Lowest byte_count + ack_ratio → ultra-low event-driven traffic"),
]
for i, (color, device, text) in enumerate(interpretations):
    col = i % 2
    row = i // 2
    x = Inches(0.6) + col * Inches(6.3)
    y = Inches(6.15) + row * Inches(0.42)
    txbox(sl_b, f"{device}:", x, y, Inches(1.4), Inches(0.35),
          font_size=Pt(10), bold=True, color=color)
    txbox(sl_b, text, x + Inches(1.4), y, Inches(4.8), Inches(0.35),
          font_size=Pt(10), color=C_DIM)

txbox(sl_b,
      "Brighter cell = that feature is MORE influential for identifying that device type",
      Inches(0.6), Inches(7.05), Inches(12.0), Inches(0.3),
      font_size=Pt(10), italic=True, color=C_DIM, align=PP_ALIGN.CENTER)


# ── Save ─────────────────────────────────────────────────────────────────────
prs.save(PPT_PATH)
print(f"Saved: {PPT_PATH}")
print(f"Total slides now: {len(prs.slides)}")
print()
for i, slide in enumerate(prs.slides):
    n_shapes = len(slide.shapes)
    print(f"  Slide {i+1:2d} — shapes: {n_shapes}")
