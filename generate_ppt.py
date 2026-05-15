"""
Generate IoT Device Fingerprinting Presentation
M.Tech Cyber Forensics | NIELIT Srinagar
Student: Md Ghufran Alam | NDU202400038
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os
import io
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np

OUT_FILE = "IoT_Fingerprinting_Presentation_FINAL.pptx"

# Color palette
BG       = RGBColor(0x0D, 0x1B, 0x2A)   # dark navy
ACCENT   = RGBColor(0x00, 0xD4, 0xFF)   # cyan
GREEN    = RGBColor(0x00, 0xFF, 0x88)   # green
ORANGE   = RGBColor(0xFF, 0x6B, 0x35)   # orange
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY    = RGBColor(0xCC, 0xCC, 0xCC)
YELLOW   = RGBColor(0xFF, 0xD7, 0x00)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # blank layout

# ─────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────
def add_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, text, left, top, width, height,
        size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        bg=None, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    if bg:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg
    return txBox

def rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def hline(slide, top, color=ACCENT, thickness=3):
    from pptx.util import Pt as PtU
    line = slide.shapes.add_connector(
        1, Inches(0.3), Inches(top), Inches(13.0), Inches(top)
    )
    line.line.color.rgb = color
    line.line.width = Pt(thickness)

def bullet_box(slide, items, left, top, width, height, size=16, color=WHITE, indent="  • "):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = indent + item
        run.font.size  = Pt(size)
        run.font.color.rgb = color

# ─────────────────────────────────────────────
# SLIDE 1 — Title
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
rect(s, 0, 0, 13.33, 0.08, ACCENT)
rect(s, 0, 7.42, 13.33, 0.08, ACCENT)

# NIELIT logo — left and right corners
LOGO = "nielit_logo.jpg"
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, Inches(0.2),  Inches(0.12), Inches(1.9), Inches(0.8))
    s.shapes.add_picture(LOGO, Inches(11.23), Inches(0.12), Inches(1.9), Inches(0.8))

# Institute name between logos
box(s, "NATIONAL INSTITUTE OF ELECTRONICS AND INFORMATION TECHNOLOGY, SRINAGAR",
    2.2, 0.18, 8.9, 0.45, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
box(s, "Ministry of Electronics & Information Technology, Govt. of India",
    2.2, 0.60, 8.9, 0.35, size=10, bold=False, color=LGRAY, align=PP_ALIGN.CENTER)

box(s, "IoT Device Fingerprinting &", 0.5, 1.0, 12.3, 1.2,
    size=36, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
box(s, "Anomaly Detection Framework", 0.5, 1.9, 12.3, 1.0,
    size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

hline(s, 3.1)

box(s, "Real-time IoT Security using Machine Learning & SHAP Explainability",
    0.5, 3.2, 12.3, 0.7, size=18, color=LGRAY, align=PP_ALIGN.CENTER)

box(s, "Md Ghufran Alam  |  Roll No. NDU202400038",
    0.5, 4.3, 12.3, 0.5, size=17, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
box(s, "M.Tech Cyber Forensics  |  NIELIT Srinagar  |  2024–2026",
    0.5, 4.85, 12.3, 0.5, size=15, color=LGRAY, align=PP_ALIGN.CENTER)
box(s, "Supervised by:  Dr. Syed Nisar Hussain Bukhari",
    0.5, 5.38, 12.3, 0.45, size=15, color=LGRAY, align=PP_ALIGN.CENTER)

rect(s, 2.5, 5.98, 8.33, 0.6, RGBColor(0x00, 0x40, 0x60))
box(s, "LIVE:  https://mdghufran-iot-fingerprinting.hf.space",
    2.5, 6.0, 8.33, 0.55, size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 2 — Problem Statement
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
rect(s, 0, 0, 13.33, 0.08, ACCENT)

box(s, "Problem Statement", 0.4, 0.15, 12.5, 0.65,
    size=28, bold=True, color=ACCENT)
hline(s, 0.9)

problems = [
    "600,000+ IoT devices compromised by Mirai botnet (2016) — took down Twitter, Netflix, Reddit",
    "Smart cameras, thermostats, bulbs cannot run antivirus software",
    "Traditional firewalls don't know WHAT device they're protecting",
    "No visibility into per-device normal behaviour baseline",
    "Botnet traffic engineered to BLEND with normal IoT patterns",
]
bullet_box(s, problems, 0.5, 1.0, 8.5, 4.0, size=17, color=WHITE)

# Right panel — impact box
rect(s, 9.3, 1.0, 3.7, 5.5, RGBColor(0x1A, 0x00, 0x00))
box(s, "Real-World Impact", 9.4, 1.05, 3.5, 0.5,
    size=14, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
impacts = ["620 Gbps DDoS peak", "12,000+ domains hit", "Twitter offline", "Netflix offline", "Reddit offline", "GitHub offline"]
bullet_box(s, impacts, 9.4, 1.6, 3.4, 3.5, size=14, color=ORANGE, indent="  ⚠ ")

box(s, "This Framework Answers:", 0.5, 5.2, 8.5, 0.45,
    size=16, bold=True, color=YELLOW)
answers = [
    'Which device is this on my network?  →  Device Fingerprinting',
    'Is it behaving normally right now?   →  Anomaly Detection',
    'Why did the model flag this?          →  SHAP Explainability',
]
bullet_box(s, answers, 0.5, 5.65, 12.5, 1.5, size=14, color=GREEN, indent="  ✓ ")

# ─────────────────────────────────────────────
# SLIDE 3 — Architecture
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
rect(s, 0, 0, 13.33, 0.08, ACCENT)

box(s, "System Architecture", 0.4, 0.15, 12.5, 0.65,
    size=28, bold=True, color=ACCENT)
hline(s, 0.9)

# Flow diagram boxes
def arch_box(sl, text, l, t, w=2.2, h=0.6, bg=RGBColor(0x00,0x40,0x60), fc=WHITE, sz=13):
    r = rect(sl, l, t, w, h, bg)
    box(sl, text, l, t+0.05, w, h-0.1, size=sz, bold=True, color=fc, align=PP_ALIGN.CENTER)

arch_box(s, "Network Traffic\n(37 Features)", 5.3, 1.0, 2.7, 0.8)

# Arrow down
box(s, "▼", 6.5, 1.85, 0.4, 0.4, size=20, color=ACCENT, align=PP_ALIGN.CENTER)

arch_box(s, "RobustScaler\n(Normalization)", 5.3, 2.25, 2.7, 0.7)

box(s, "▼", 6.5, 2.99, 0.4, 0.4, size=20, color=ACCENT, align=PP_ALIGN.CENTER)

# Split
box(s, "┌─────────────────────────────────┐", 1.5, 3.4, 10.0, 0.4, size=14, color=ACCENT)

arch_box(s, "Device Fingerprinting\nRandom Forest", 1.0, 3.6, 3.5, 0.8, bg=RGBColor(0x00,0x50,0x30))
arch_box(s, "Anomaly Detection\nIF + OC-SVM Ensemble", 5.2, 3.6, 3.5, 0.8, bg=RGBColor(0x50,0x20,0x00))
arch_box(s, "SHAP Explainability\nTreeExplainer", 9.4, 3.6, 3.5, 0.8, bg=RGBColor(0x30,0x00,0x50))

box(s, "▼", 2.6, 4.45, 0.4, 0.35, size=18, color=GREEN, align=PP_ALIGN.CENTER)
box(s, "▼", 6.8, 4.45, 0.4, 0.35, size=18, color=ORANGE, align=PP_ALIGN.CENTER)
box(s, "▼", 11.0, 4.45, 0.4, 0.35, size=18, color=YELLOW, align=PP_ALIGN.CENTER)

arch_box(s, "Device Label\n+ Confidence", 1.0, 4.8, 3.5, 0.7, bg=RGBColor(0x00,0x35,0x20))
arch_box(s, "Anomaly Score\n+ Severity", 5.2, 4.8, 3.5, 0.7, bg=RGBColor(0x35,0x15,0x00))
arch_box(s, "Top-10 Feature\nContributions", 9.4, 4.8, 3.5, 0.7, bg=RGBColor(0x20,0x00,0x35))

box(s, "▼", 6.5, 5.55, 0.4, 0.35, size=18, color=ACCENT, align=PP_ALIGN.CENTER)

arch_box(s, "Alert Manager  →  FastAPI REST API  +  Live Plotly.js Dashboard", 1.0, 5.9, 11.5, 0.7,
         bg=RGBColor(0x00,0x20,0x40), sz=15)

# ─────────────────────────────────────────────
# SLIDE 4 — Dataset & Features
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
rect(s, 0, 0, 13.33, 0.08, ACCENT)

box(s, "Dataset & Feature Engineering", 0.4, 0.15, 12.5, 0.65,
    size=28, bold=True, color=ACCENT)
hline(s, 0.9)

# Left — dataset info
box(s, "N-BaIoT Dataset  (UCI ML Repository #442)", 0.4, 1.0, 6.5, 0.5,
    size=18, bold=True, color=YELLOW)
ds_items = [
    "9 real physical IoT devices",
    "Mirai botnet: ack, scan, syn, udp, udpplain attacks",
    "BASHLITE/Gafgyt: combo, junk, scan, tcp, udp",
    "Normal traffic baseline per device",
    "115 original features → mapped to 37",
]
bullet_box(s, ds_items, 0.4, 1.55, 6.5, 3.0, size=15)

box(s, "Device Mapping:", 0.4, 4.5, 6.5, 0.4, size=15, bold=True, color=ACCENT)
devices = [
    "Danmini Doorbell      →  smart_doorbell",
    "Ecobee Thermostat     →  smart_thermostat",
    "Security Camera       →  smart_camera",
    "Baby Monitor          →  smart_speaker",
    "Samsung Webcam        →  smart_tv",
    "Philips Camera        →  smart_bulb",
    "SimpleHome Camera     →  smart_plug",
    "Ennio Doorbell        →  motion_sensor",
]
bullet_box(s, devices, 0.4, 4.9, 6.5, 2.5, size=13, color=LGRAY)

# Right — 37 features table
rect(s, 7.0, 1.0, 6.0, 6.2, RGBColor(0x05, 0x15, 0x25))
box(s, "37 Network Flow Features", 7.1, 1.05, 5.8, 0.45,
    size=15, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

feature_cats = [
    ("Temporal  (5)", "flow_duration, mean_iat, std_iat, min_iat, max_iat"),
    ("Volume  (4)", "packet_count, byte_count, packet_rate, byte_rate"),
    ("Packet Size  (4)", "mean_pkt_size, std_pkt_size, min_pkt_size, max_pkt_size"),
    ("Protocol  (6)", "tcp_ratio, udp_ratio, syn_ratio, fin_ratio, rst_ratio, ack_ratio"),
    ("App Layer  (8)", "is_https, is_mqtt, is_coap, is_mdns, is_ntp, dns_queries, well_known_ports, is_encrypted"),
    ("Direction  (3)", "upload_bytes, download_bytes, upload_download_ratio"),
    ("Destination  (7)", "unique_dest_ports, unique_dest_ips, port_entropy, ip_entropy, well_known_ports_count, mean_dest_port, std_dest_port"),
]
y = 1.6
for cat, feats in feature_cats:
    box(s, cat, 7.1, y, 2.0, 0.38, size=12, bold=True, color=GREEN)
    box(s, feats, 9.1, y, 3.8, 0.5, size=10, color=LGRAY)
    y += 0.72

# ─────────────────────────────────────────────
# SLIDE 5 — ML Models
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
rect(s, 0, 0, 13.33, 0.08, ACCENT)

box(s, "Machine Learning Models", 0.4, 0.15, 12.5, 0.65,
    size=28, bold=True, color=ACCENT)
hline(s, 0.9)

# Fingerprinting models
box(s, "MODULE 1 — Device Fingerprinting", 0.4, 1.0, 6.2, 0.5,
    size=18, bold=True, color=GREEN)

fp_models = [
    ("Random Forest (Primary)", "100 trees, Gini, balanced, 100% acc"),
    ("Gradient Boosting", "100 trees, lr=0.1, max_depth=5"),
    ("SVM (RBF kernel)", "C=1.0, gamma=scale, 94.89% acc"),
    ("Voting Ensemble", "RF + GB + SVM soft vote, 100% acc"),
]
y = 1.6
for name, detail in fp_models:
    rect(s, 0.4, y, 6.2, 0.55, RGBColor(0x05, 0x25, 0x15))
    box(s, name, 0.5, y+0.04, 3.0, 0.45, size=13, bold=True, color=WHITE)
    box(s, detail, 3.5, y+0.06, 3.0, 0.42, size=12, color=LGRAY)
    y += 0.65

box(s, "Training: SMOTE balancing + RobustScaler", 0.4, 4.3, 6.2, 0.4,
    size=13, color=YELLOW)
box(s, "Confidence threshold: 60% (below = Unknown)", 0.4, 4.7, 6.2, 0.4,
    size=13, color=LGRAY)

# Anomaly detection models
box(s, "MODULE 2 — Anomaly Detection", 7.0, 1.0, 6.0, 0.5,
    size=18, bold=True, color=ORANGE)

ad_text = [
    "Per-device models — each device type gets its OWN detector",
    "Trained on NORMAL traffic ONLY (unsupervised)",
    "No labeled attack data required",
    "",
    "Ensemble per device:",
    "  • Isolation Forest  (contamination=0.1)",
    "  • One-Class SVM  (kernel=rbf, nu=0.1)",
    "  • Final score = weighted average",
    "",
    "Threshold: 0.75",
    "  < 0.75  →  Normal",
    "  0.50–0.65  →  LOW alert",
    "  0.65–0.80  →  MEDIUM alert",
    "  0.80–0.90  →  HIGH alert",
    "  0.90+       →  CRITICAL alert",
]
y = 1.6
for line in ad_text:
    color = ORANGE if line.startswith("Ensemble") or line.startswith("Threshold") else LGRAY
    if line.startswith("  •"):
        color = WHITE
    if "0.90+" in line or "CRITICAL" in line:
        color = RGBColor(0xFF, 0x44, 0x44)
    elif "HIGH" in line:
        color = ORANGE
    elif "MEDIUM" in line:
        color = YELLOW
    elif "LOW" in line:
        color = GREEN
    box(s, line, 7.0, y, 5.8, 0.35, size=13, color=color)
    y += 0.33

# ─────────────────────────────────────────────
# SLIDE 6 — Results: Fingerprinting
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
rect(s, 0, 0, 13.33, 0.08, ACCENT)

box(s, "Results — Device Fingerprinting", 0.4, 0.15, 12.5, 0.65,
    size=28, bold=True, color=ACCENT)
hline(s, 0.9)

# Model comparison table
headers = ["Model", "Test Accuracy", "ROC-AUC", "F1 (Macro)"]
rows = [
    ["Random Forest  ★", "100.00%", "1.0000", "1.00"],
    ["Gradient Boosting", "100.00%", "1.0000", "1.00"],
    ["Voting Ensemble", "100.00%", "1.0000", "1.00"],
    ["SVM (RBF)", "94.89%", "0.9967", "0.95"],
]
col_x = [0.4, 4.5, 7.5, 10.5]
col_w = [4.0, 2.8, 2.8, 2.5]

rect(s, 0.4, 1.0, 12.5, 0.5, RGBColor(0x00, 0x40, 0x60))
for i, h in enumerate(headers):
    box(s, h, col_x[i], 1.0, col_w[i], 0.5, size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

for ri, row in enumerate(rows):
    bg_color = RGBColor(0x00, 0x30, 0x10) if ri == 0 else RGBColor(0x05, 0x10, 0x20)
    rect(s, 0.4, 1.55 + ri*0.5, 12.5, 0.48, bg_color)
    fc = GREEN if ri == 0 else WHITE
    for ci, cell in enumerate(row):
        box(s, cell, col_x[ci], 1.57 + ri*0.5, col_w[ci], 0.44,
            size=14, bold=(ri==0), color=fc, align=PP_ALIGN.CENTER)

box(s, "Per-Device Classification Report:", 0.4, 3.65, 6.2, 0.4,
    size=16, bold=True, color=YELLOW)

per_device = [
    ("smart_camera", "100%", "1.000", "1.000"),
    ("smart_thermostat", "100%", "1.000", "1.000"),
    ("smart_tv", "100%", "1.000", "1.000"),
    ("smart_bulb", "100%", "1.000", "1.000"),
    ("motion_sensor", "100%", "1.000", "1.000"),
    ("smart_plug", "100%", "1.000", "1.000"),
    ("smart_speaker", "100%", "1.000", "1.000"),
    ("smart_doorbell", "100%", "1.000", "1.000"),
]
y = 4.1
for dev, acc, prec, rec in per_device:
    box(s, f"{dev:<22}  Acc: {acc}   Precision: {prec}   Recall: {rec}",
        0.5, y, 8.0, 0.35, size=13, color=LGRAY)
    y += 0.37

# Right — key insight
rect(s, 9.2, 3.65, 3.8, 3.7, RGBColor(0x00, 0x25, 0x10))
box(s, "Key Insight", 9.3, 3.7, 3.6, 0.45, size=15, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
insights = [
    "Pure network traffic — no agent, no device access",
    "37 features enough for 100% separation",
    "Protocol ratios (MQTT, CoAP, HTTPS) are the most discriminative",
    "Works across all 8 device categories without overlap",
]
bullet_box(s, insights, 9.3, 4.2, 3.6, 3.0, size=13, color=WHITE)

# ─────────────────────────────────────────────
# SLIDE 7 — Results: Anomaly Detection
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
rect(s, 0, 0, 13.33, 0.08, ACCENT)

box(s, "Results — Anomaly Detection", 0.4, 0.15, 12.5, 0.65,
    size=28, bold=True, color=ACCENT)
hline(s, 0.9)

box(s, "Per-Device Anomaly Detection Performance (AUC-ROC)", 0.4, 1.0, 8.0, 0.45,
    size=17, bold=True, color=YELLOW)

ad_results = [
    ("smart_doorbell",   0.8773, 0.4811),
    ("smart_thermostat", 0.8590, 0.3963),
    ("motion_sensor",    0.7946, 0.2206),
    ("smart_tv",         0.7732, 0.2415),
    ("smart_plug",       0.7341, 0.2990),
    ("smart_camera",     0.7125, 0.2610),
]

y = 1.55
for dev, auc, f1 in ad_results:
    bar_w = auc * 5.0
    rect(s, 0.4, y, 8.5, 0.52, RGBColor(0x08, 0x12, 0x20))
    box(s, f"{dev:<22}", 0.5, y+0.08, 2.8, 0.36, size=13, color=WHITE)
    rect(s, 3.3, y+0.12, bar_w, 0.28, ORANGE if auc >= 0.80 else ACCENT)
    box(s, f"AUC: {auc:.4f}   F1: {f1:.4f}", 8.6, y+0.08, 4.4, 0.36, size=13, color=LGRAY)
    y += 0.6

box(s, "Why F1 is lower than AUC-ROC:", 0.4, 5.35, 8.0, 0.4,
    size=15, bold=True, color=ORANGE)
f1_note = [
    "Mirai & BASHLITE engineered to BLEND with normal IoT traffic",
    "Unsupervised model never saw attack data during training",
    "AUC-ROC is threshold-independent — correct metric for anomaly detection",
    "F1 depends on threshold — can be tuned per deployment",
]
bullet_box(s, f1_note, 0.4, 5.78, 8.0, 1.6, size=14, color=LGRAY)

# Right panel
rect(s, 9.0, 1.0, 4.0, 6.3, RGBColor(0x15, 0x05, 0x00))
box(s, "Live Demo Results", 9.1, 1.05, 3.8, 0.45, size=15, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
demo_stats = [
    "Normal traffic → score 0.00",
    "Data exfiltration → 0.81 HIGH",
    "Port scan attack → 0.77 MED",
    "Combined /analyze → 0.81 HIGH",
    "",
    "Zero false positives",
    "on normal traffic",
    "",
    "100% detection rate",
    "on injected attacks",
    "",
    "Threshold 0.75",
    "cleanly separates",
    "normal vs attack",
]
bullet_box(s, demo_stats, 9.1, 1.6, 3.8, 5.5, size=13, color=WHITE)

# ─────────────────────────────────────────────
# SHAP bar chart (generated in-memory)
# ─────────────────────────────────────────────
_shap_features = [
    ("max_pkt_size",          +0.04),
    ("is_mqtt",               +0.07),
    ("syn_ratio",             +0.08),
    ("well_known_port_ratio", -0.05),
    ("is_encrypted",          -0.06),
    ("mean_dest_port",        -0.09),
    ("is_https",              -0.10),
    ("ack_ratio",             -0.12),
    ("udp_ratio",             -0.15),
    ("tcp_ratio",             -0.18),
]
_feat_names = [f[0] for f in _shap_features]
_shap_vals  = [f[1] for f in _shap_features]
_colors     = ['#00FF88' if v > 0 else '#FF5555' for v in _shap_vals]

fig, ax = plt.subplots(figsize=(5.5, 3.6), facecolor='#0D1B2A')
ax.set_facecolor('#0D1B2A')
bars = ax.barh(_feat_names, _shap_vals, color=_colors, height=0.6, edgecolor='none')
ax.axvline(0, color='#AAAAAA', linewidth=0.8)
ax.set_xlabel('SHAP Value (impact on prediction)', color='#CCCCCC', fontsize=9)
ax.set_title('Unknown Device — 46% confidence', color='#00D4FF', fontsize=10, fontweight='bold')
ax.tick_params(colors='#CCCCCC', labelsize=8)
for spine in ax.spines.values():
    spine.set_edgecolor('#333333')
ax.xaxis.label.set_color('#CCCCCC')
green_patch = mpatches.Patch(color='#00FF88', label='Toward (supports class)')
red_patch   = mpatches.Patch(color='#FF5555', label='Away (reduces confidence)')
ax.legend(handles=[green_patch, red_patch], facecolor='#0D1B2A',
          labelcolor='#CCCCCC', fontsize=7, loc='lower right')
plt.tight_layout(pad=0.4)

_shap_buf = io.BytesIO()
fig.savefig(_shap_buf, format='png', dpi=150, facecolor='#0D1B2A')
plt.close(fig)
_shap_buf.seek(0)

# ─────────────────────────────────────────────
# SLIDE 8 — SHAP Explainability
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
rect(s, 0, 0, 13.33, 0.08, ACCENT)

box(s, "SHAP Explainability (XAI)", 0.4, 0.15, 12.5, 0.65,
    size=28, bold=True, color=ACCENT)
hline(s, 0.9)

box(s, "Why does explainability matter in security?", 0.4, 1.05, 9.0, 0.45,
    size=17, bold=True, color=YELLOW)
xai_why = [
    "A black-box alert ('anomaly detected') is not actionable for a SOC analyst",
    "SHAP shows WHICH features triggered the alert — analyst can verify manually",
    "Regulatory compliance (GDPR Article 22) requires explainable automated decisions",
    "Reduces false positive investigation time from minutes to seconds",
]
bullet_box(s, xai_why, 0.4, 1.55, 9.0, 2.0, size=15, color=WHITE)

hline(s, 3.65, color=RGBColor(0x33,0x33,0x33), thickness=1)

box(s, "Top-10 SHAP Features — Unknown Device (46% confidence)", 0.4, 3.75, 9.0, 0.45,
    size=16, bold=True, color=ACCENT)

shap_features = [
    ("tcp_ratio",             "Red  (Away)", "High TCP usage inconsistent with device profile"),
    ("udp_ratio",             "Red  (Away)", "UDP ratio too high for predicted class"),
    ("ack_ratio",             "Red  (Away)", "ACK pattern deviates from baseline"),
    ("is_https",              "Red  (Away)", "HTTPS presence reduces confidence"),
    ("mean_dest_port",        "Red  (Away)", "Port profile does not match any device"),
    ("syn_ratio",             "Green (Toward)", "SYN ratio partially supports class"),
    ("is_mqtt",               "Green (Toward)", "MQTT traffic partially consistent"),
    ("is_encrypted",          "Red  (Away)", "Encryption profile deviates"),
    ("well_known_port_ratio", "Red  (Away)", "Port ratio outside expected range"),
    ("max_pkt_size",          "Green (Toward)", "Packet size aligns with device"),
]

y = 4.25
for feat, direction, interp in shap_features:
    color = GREEN if "Green" in direction else RGBColor(0xFF, 0x55, 0x55)
    box(s, f"{feat:<25}  {direction:<16}  {interp}", 0.4, y, 12.5, 0.32,
        size=12, color=color)
    y += 0.31

# SHAP bar chart — right side
pic = s.shapes.add_picture(_shap_buf, Inches(8.6), Inches(1.0), Inches(4.5), Inches(3.2))

# Implementation box below chart
rect(s, 8.6, 4.3, 4.5, 2.8, RGBColor(0x05, 0x05, 0x20))
box(s, "SHAP Implementation", 8.7, 4.35, 4.3, 0.4, size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
shap_impl = [
    "TreeExplainer (faster than KernelExplainer)",
    "Runs on Random Forest model",
    "POST /explain endpoint",
    "Live panel in monitoring dashboard",
    "Updates every prediction",
    "Sub-30ms response time",
]
bullet_box(s, shap_impl, 8.7, 4.78, 4.3, 2.2, size=12, color=LGRAY)

# ─────────────────────────────────────────────
# SLIDE 9 — Live Deployment
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
rect(s, 0, 0, 13.33, 0.08, ACCENT)

box(s, "Live Deployment", 0.4, 0.15, 12.5, 0.65,
    size=28, bold=True, color=ACCENT)
hline(s, 0.9)

# Platform info
rect(s, 0.4, 1.0, 5.5, 1.2, RGBColor(0x00, 0x20, 0x35))
box(s, "Platform: HuggingFace Spaces (Docker)", 0.5, 1.05, 5.3, 0.45, size=14, bold=True, color=YELLOW)
box(s, "SDK: Docker  |  Visibility: Public  |  Free tier", 0.5, 1.5, 5.3, 0.4, size=13, color=LGRAY)

rect(s, 6.4, 1.0, 6.5, 1.2, RGBColor(0x00, 0x20, 0x35))
box(s, "Single Port Deployment (7860)", 6.5, 1.05, 6.3, 0.45, size=14, bold=True, color=YELLOW)
box(s, "FastAPI serves API + Dashboard on single port 7860", 6.5, 1.5, 6.3, 0.4, size=13, color=LGRAY)

# URL list
urls = [
    ("Live Dashboard",   "https://mdghufran-iot-fingerprinting.hf.space",         GREEN),
    ("API Docs",         "https://mdghufran-iot-fingerprinting.hf.space/docs",     ACCENT),
    ("Status Page",      "https://mdghufran-iot-fingerprinting.hf.space/status",   ACCENT),
    ("HuggingFace",      "https://huggingface.co/spaces/mdghufran/iot-fingerprinting", YELLOW),
    ("GitHub",           "https://github.com/Ghufran2002/iot-fingerprinting-framework", LGRAY),
]
y = 2.4
for label, url, color in urls:
    rect(s, 0.4, y, 12.5, 0.52, RGBColor(0x05, 0x15, 0x25))
    box(s, f"{label:<18}", 0.5, y+0.07, 2.2, 0.38, size=14, bold=True, color=WHITE)
    box(s, url, 2.7, y+0.07, 10.0, 0.38, size=14, color=color)
    y += 0.6

# API endpoints
box(s, "API Endpoints:", 0.4, 5.5, 5.0, 0.4, size=15, bold=True, color=ACCENT)
endpoints = [
    "GET  /health    GET  /devices    GET  /metrics",
    "POST /fingerprint    POST /anomaly/score    POST /analyze    POST /explain",
]
bullet_box(s, endpoints, 0.4, 5.92, 9.0, 1.3, size=13, color=LGRAY, indent="  ")

rect(s, 9.5, 5.5, 3.5, 1.7, RGBColor(0x00, 0x30, 0x10))
box(s, "Performance", 9.6, 5.55, 3.3, 0.4, size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
perf = ["API latency: ~40 ms", "Status: Live & Healthy", "Models: RF + IF + OC-SVM", "Accuracy: 100% (8 devices)"]
bullet_box(s, perf, 9.6, 5.98, 3.3, 1.1, size=13, color=WHITE)

# ─────────────────────────────────────────────
# SLIDE 10 — Demo: Fingerprinting
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
rect(s, 0, 0, 13.33, 0.08, ACCENT)

box(s, "Live Demo — Device Fingerprinting", 0.4, 0.15, 12.5, 0.65,
    size=28, bold=True, color=ACCENT)
hline(s, 0.9)

box(s, "POST /fingerprint  — 8 devices identified from network traffic alone", 0.4, 1.0, 12.5, 0.45,
    size=16, color=LGRAY)

devices_demo = [
    ("192.168.1.45", "smart_camera",     "99%", "HTTPS, 400 pkt/s, large packets, port 443"),
    ("192.168.1.22", "smart_thermostat", "98%", "MQTT, 2 pkt/s, 80-byte packets, port 1883"),
    ("192.168.1.10", "smart_tv",         "97%", "Streaming, 695 MB download, 18 IPs"),
    ("192.168.1.77", "smart_bulb",       "96%", "CoAP, 4 packets, 120 bytes, port 5683"),
    ("192.168.1.99", "motion_sensor",    "95%", "MQTT+CoAP, 3 packets, event-driven"),
    ("192.168.1.55", "smart_plug",       "97%", "MQTT, 6 packets, port 8883"),
    ("192.168.1.33", "smart_speaker",    "98%", "HTTPS, voice streaming, 5 cloud IPs"),
    ("192.168.1.88", "smart_doorbell",   "96%", "Camera+audio burst, 2.5 MB upload, HTTPS"),
]

y = 1.55
for ip, dev, conf, sig in devices_demo:
    rect(s, 0.4, y, 12.5, 0.48, RGBColor(0x05, 0x15, 0x25))
    box(s, ip, 0.5, y+0.07, 1.8, 0.34, size=12, color=LGRAY)
    box(s, dev, 2.3, y+0.07, 2.5, 0.34, size=13, bold=True, color=GREEN)
    rect(s, 4.8, y+0.1, 0.9, 0.28, RGBColor(0x00, 0x40, 0x20))
    box(s, conf, 4.8, y+0.07, 0.9, 0.34, size=12, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
    box(s, sig, 5.8, y+0.07, 7.0, 0.34, size=12, color=LGRAY)
    y += 0.54

box(s, "No hardware access. No credentials. No agent installed on any device.", 0.4, 6.45, 12.5, 0.45,
    size=15, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 11 — Demo: Anomaly Detection
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
rect(s, 0, 0, 13.33, 0.08, ACCENT)

box(s, "Live Demo — Anomaly Detection", 0.4, 0.15, 12.5, 0.65,
    size=28, bold=True, color=ACCENT)
hline(s, 0.9)

box(s, "POST /anomaly/score  +  POST /analyze", 0.4, 1.0, 12.5, 0.4, size=16, color=LGRAY)

tests = [
    ("TEST 1", "Normal Camera Traffic",        "0.0000", "NORMAL",   RGBColor(0x00,0x40,0x20), GREEN,  "400 pkt/s, 28MB upload, ratio=2 → System: NORMAL"),
    ("TEST 2", "Data Exfiltration (Mirai)",    "0.8060", "HIGH",     RGBColor(0x40,0x18,0x00), ORANGE, "Camera uploading 1.68 GB, ratio=120 (normally 2) → ALERT"),
    ("TEST 3", "Port Scan / Reconnaissance",   "0.7745", "MEDIUM",   RGBColor(0x30,0x25,0x00), YELLOW, "Thermostat scanning 1200 ports × 250 IPs, entropy=9.5 → ALERT"),
    ("TEST 4", "/analyze Combined Call",       "0.8060", "HIGH",     RGBColor(0x40,0x18,0x00), ORANGE, "Auto-detected camera → exfiltration. One API call does both"),
]

y = 1.5
for label, name, score, sev, bg, sc, desc in tests:
    rect(s, 0.4, y, 12.5, 1.0, bg)
    box(s, label, 0.5, y+0.05, 1.0, 0.4, size=13, bold=True, color=WHITE)
    box(s, name, 1.5, y+0.05, 4.0, 0.4, size=14, bold=True, color=WHITE)
    box(s, f"Score: {score}", 5.5, y+0.05, 2.0, 0.4, size=14, bold=True, color=sc)
    rect(s, 7.5, y+0.08, 1.5, 0.3, RGBColor(0x10,0x10,0x10))
    box(s, sev, 7.5, y+0.05, 1.5, 0.35, size=13, bold=True, color=sc, align=PP_ALIGN.CENTER)
    box(s, desc, 0.5, y+0.55, 12.0, 0.38, size=12, color=LGRAY)
    y += 1.1

rect(s, 0.4, 5.9, 12.5, 0.55, RGBColor(0x00, 0x25, 0x00))
box(s, "Zero false positives on normal traffic.  All 3 attack types detected with correct severity.",
    0.5, 5.93, 12.3, 0.45, size=15, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 12 — Conclusion
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
rect(s, 0, 0, 13.33, 0.08, ACCENT)

box(s, "Conclusion & Contributions", 0.4, 0.15, 12.5, 0.65,
    size=28, bold=True, color=ACCENT)
hline(s, 0.9)

contributions = [
    ("Fingerprinting", "100% accuracy on 8 IoT device types using only network traffic — no device access required"),
    ("Anomaly Detection", "Per-device unsupervised ensemble (IF + OC-SVM) detects Mirai/BASHLITE with AUC up to 0.877"),
    ("Explainability", "SHAP TreeExplainer provides real-time top-10 feature contributions for every prediction"),
    ("Production Deploy", "Live on HuggingFace Spaces — publicly accessible, sub-40ms latency, 2.1h+ uptime"),
    ("REST API", "8 FastAPI endpoints including combined /analyze, /explain, /alerts/recent"),
    ("Dashboard", "Live Plotly.js monitoring dashboard with anomaly timeline, severity chart, real-time alert feed"),
    ("Dataset", "Trained on real N-BaIoT dataset (UCI #442) with real Mirai & BASHLITE attack traffic"),
    ("Open Source", "Full code available on GitHub — reproducible, extensible, MIT licensed"),
]

y = 1.05
for title, desc in contributions:
    rect(s, 0.4, y, 12.5, 0.52, RGBColor(0x05, 0x12, 0x22))
    box(s, f"✓  {title}", 0.5, y+0.07, 2.8, 0.38, size=13, bold=True, color=GREEN)
    box(s, desc, 3.3, y+0.07, 9.5, 0.38, size=13, color=WHITE)
    y += 0.58

rect(s, 0.4, 5.85, 12.5, 0.75, RGBColor(0x00, 0x18, 0x30))
box(s, "LIVE AT:  https://mdghufran-iot-fingerprinting.hf.space",
    0.5, 5.88, 12.3, 0.65, size=18, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 13 — Thank You
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s)
rect(s, 0, 0, 13.33, 0.08, ACCENT)
rect(s, 0, 7.42, 13.33, 0.08, ACCENT)

box(s, "Thank You", 0.5, 1.6, 12.3, 1.2,
    size=52, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
hline(s, 3.0)
box(s, "Questions & Discussion", 0.5, 3.15, 12.3, 0.7,
    size=24, color=WHITE, align=PP_ALIGN.CENTER)

box(s, "Md Ghufran Alam  |  NDU202400038",
    0.5, 4.1, 12.3, 0.55, size=20, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
box(s, "M.Tech Cyber Forensics  |  NIELIT Srinagar  |  2024–2026",
    0.5, 4.7, 12.3, 0.5, size=16, color=LGRAY, align=PP_ALIGN.CENTER)

rect(s, 1.5, 5.5, 10.3, 1.5, RGBColor(0x00, 0x20, 0x35))
box(s, "Live System:  https://mdghufran-iot-fingerprinting.hf.space",
    1.6, 5.55, 10.1, 0.5, size=16, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
box(s, "GitHub:  https://github.com/Ghufran2002/iot-fingerprinting-framework",
    1.6, 6.05, 10.1, 0.45, size=14, color=LGRAY, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
prs.save(OUT_FILE)
print(f"\nPresentation saved: {OUT_FILE}")
print(f"Slides: {len(prs.slides)}")
print("Open the file to view the presentation.")
